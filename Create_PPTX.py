from pathlib import Path
from pptx import Presentation
from datetime import datetime
from pptx.dml.color import RGBColor
import copy


AUSWERTUNG = Path.cwd() / "Auswertung"
EXCEL_GRAPHS = Path.cwd() / "Graphs"
PPTX_TEMPLATE = Path("PPTX_files/DATUM_Messdatenauswertung_MASCHINE_Cone_V2.pptx")
TESTs = ["FPPSN2", "FPPSN6", "DPPSN0", "FPPSN4"]
TEST_STRING = "FPPSN2, FPPSN6, DPPSN0, FPPSN4"
Project = "HD11EA"

FOOTER_TEXT = f'{TEST_STRING} | Conceptone | {datetime.now().strftime("%d-%m-%Y")}'


def check_shape_overlap(shapes, left, top):
    """
    Check if any shape in the given collection overlaps with the specified position.
    """
    for shape in shapes:
        if shape.left == left and shape.top == top:
            return True
    return False


def CreateNewSetOfSlides(base_pptx, slide_index, start_index):

    test_name_layout_number = 1
    cycle_layout_number = 5
    einfall_layout_number = 6
    dauer_layout_number = 7

    layout_slide = {
        0: test_name_layout_number,
        1: cycle_layout_number,
        2: einfall_layout_number,
        3: dauer_layout_number,
        4: cycle_layout_number,
        5: einfall_layout_number,
        6: dauer_layout_number
    }

    for add_index, layout in layout_slide.items():
        source = base_pptx.slides[start_index + add_index]
        destination = base_pptx.slides.add_slide(base_pptx.slide_layouts[layout])

        for shp in source.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            if not check_shape_overlap(destination.shapes, shp.left, shp.top):
                destination.shapes._spTree.insert_element_before(newel, 'p:extLst')

        xml_slides = base_pptx.slides._sldIdLst
        old_index = -1 # Slide is Added to End
        new_index = slide_index + add_index
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    return new_index + 1


def SetFooterText(shape):
    shape.text = FOOTER_TEXT
    paragraph = shape.text_frame.paragraphs[0]
    font = paragraph.font
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)


def CycleImageSlide(current_slide_shapes, test_name, variant):
    for placeholder_index, shape in enumerate(current_slide_shapes):
        if placeholder_index == 0:
            shape.text = f"100 / 100 Zyklen, Ein-Lückenrad (100 x {variant}): Andrehen auf n_ab =18 (n_an=173rpm) bis PS fängt"
        if placeholder_index == 1:
            SetFooterText(shape)
        if placeholder_index == 3:
            shape.insert_picture(f"{AUSWERTUNG / f'{test_name}_{variant}.png'}")


def DauerSlide(current_slide_shapes, test_name, variant):
    for placeholder_index, shape in enumerate(current_slide_shapes):
        if placeholder_index == 0:
            SetFooterText(shape)
        if placeholder_index == 2:
            shape.insert_picture(f"{EXCEL_GRAPHS / f'DauerBisEinfall_{test_name}_{variant}.png'}")


def EinfallSlide(current_slide_shapes, test_name, variant):
    for placeholder_index, shape in enumerate(current_slide_shapes):
        if placeholder_index == 0:
            SetFooterText(shape)
        if placeholder_index == 2:
            shape.insert_picture(f"{EXCEL_GRAPHS / f'Einfalldrehzahlen_{test_name}_{variant}.png'}")


def EditFirstGraphSlides(base_pptx_slides, test_name, start_index, end_index):

    for index, current_slide in enumerate(base_pptx_slides):
        if index == start_index:
            for placeholder_index, shape in enumerate(current_slide.shapes):
                if placeholder_index == 0:
                    shape.text = "01"
                if placeholder_index == 2:
                    shape.text = test_name
                if placeholder_index == 3:
                    SetFooterText(shape)

        elif index == start_index + 1:
            CycleImageSlide(current_slide.shapes, test_name, "VW")
        elif index == start_index + 2:
            EinfallSlide(current_slide.shapes, test_name, "VW")
        elif index == start_index + 3:
            DauerSlide(current_slide.shapes, test_name, "VW")
        elif index == start_index + 4:
            CycleImageSlide(current_slide.shapes, test_name, "RW")
        elif index == start_index + 5:
            EinfallSlide(current_slide.shapes, test_name, "VW")
        elif index == end_index:
            DauerSlide(current_slide.shapes, test_name, "RW")


def EditAllGraphSlides(base_pptx_slides, test_name, start_index, end_index):

    for index, current_slide in enumerate(base_pptx_slides):
        if index == start_index:
            for placeholder_index, shape in enumerate(current_slide.shapes):
                print(f"{placeholder_index} == ")

        elif index == start_index + 1:
            CycleImageSlide(current_slide.shapes, test_name, "VW")
        elif index == start_index + 2:
            EinfallSlide(current_slide.shapes, test_name, "VW")
        elif index == start_index + 3:
            DauerSlide(current_slide.shapes, test_name, "VW")
        elif index == start_index + 4:
            CycleImageSlide(current_slide.shapes, test_name, "RW")
        elif index == start_index + 5:
            EinfallSlide(current_slide.shapes, test_name, "VW")
        elif index == end_index:
            DauerSlide(current_slide.shapes, test_name, "RW")

def CreatePPTX():
    base_pptx = Presentation(PPTX_TEMPLATE)
    base_pptx_slides = base_pptx.slides
    start_index = 3
    end_index = 9
    new_slide_index = 10
    for _ in range(len(TESTs) - 1):
        new_slide_index = CreateNewSetOfSlides(base_pptx, new_slide_index, start_index)
    
    EditFirstGraphSlides(base_pptx_slides, "FPPSN2", start_index, end_index)

    
    EditAllGraphSlides(base_pptx_slides, "FPPSN6", 10, 16)

    base_pptx.save("NEW.pptx")


if __name__=="__main__":
    CreatePPTX()
